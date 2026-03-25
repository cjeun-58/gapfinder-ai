import streamlit as st
from google import genai
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from fpdf import FPDF
import io
import os
import re

# --- 1. 페이지 설정 및 데이터 초기화 ---
_ = st.set_page_config(page_title="GapFinder AI v14.2", layout="wide")

# 세션 데이터 초기화 (분석 결과 및 운영 인사이트 유지)
states = ['brand_analysis', 'brand_insight', 'comp_analysis', 'consumer_data', 'consumer_analysis', 'final_report']
for key in states:
    if key not in st.session_state:
        st.session_state[key] = [] if 'data' in key else ""

# --- 2. 사이드바 (실시간 현황 체크) ---
with st.sidebar:
    st.header("🔑 API 설정")
    gemini_key = st.text_input("1. Gemini API Key", type="password")
    serper_key = st.text_input("2. Serper API Key", type="password")
    _ = st.divider()
    menu = st.radio("전략 수립 단계", [
        "STEP 1. 자사 분석 (Thesis)", 
        "STEP 1.5. 경쟁사 Deep-Dive", 
        "STEP 2. 소비자 분석 (Antithesis)", 
        "STEP 3. 하이브리드 전략 및 PDF"
    ])
    _ = st.divider()
    st.subheader("📊 실시간 분석 현황")
    st.write(f"🏢 자사: {'✅ 완료' if st.session_state['brand_analysis'] else '❌ 미완료'}")
    st.write(f"⚔️ 경쟁사: {'✅ 완료' if st.session_state['comp_analysis'] else '❌ 미완료'}")
    st.write(f"👥 소비자: {'✅ 완료' if st.session_state['consumer_analysis'] else '❌ 미완료'}")

# --- 3. 핵심 유틸리티 함수 ---

def extract_text(files=None, urls=None):
    """파일(PDF, PPTX, XLSX)과 URL에서 텍스트를 추출합니다."""
    text = ""
    if files:
        for f in files:
            try:
                if f.name.endswith(".pdf"):
                    text += "\n".join([p.extract_text() for p in PdfReader(f).pages if p.extract_text()])
                elif f.name.endswith(".pptx"):
                    prs = Presentation(f)
                    text += "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif f.name.endswith(".xlsx"):
                    text += pd.read_excel(f).to_string()
            except Exception:
                pass
    if urls:
        for url in urls:
            if url and url.strip():
                try:
                    res = requests.get(url.strip(), headers={'User-Agent': 'Mozilla/5.0'}, timeout=5)
                    soup = BeautifulSoup(res.text, 'html.parser')
                    for s in soup(['script', 'style']): s.decompose()
                    text += f"\n[참조 내용: {url}]\n{soup.get_text()[:3000]}"
                except Exception:
                    pass
    return text

def analyze_hybrid(content, target_type, insight="", brand_context=""):
    """
    [핵심 고도화] brand_context를 경쟁사 분석 시 주입하여 주객전도를 방지합니다.
    """
    try:
        client = genai.Client(api_key=gemini_key)
        p_base = "자기소개 금지. 광고 기획자의 시각으로 리스트 형식을 사용하여 분석하세요. 마크다운 표(|)는 사용하지 마세요.\n\n"
        
        prompts = {
            "brand": f"{p_base}[Thesis] 자사 브랜드 가치와 아래 운영 인사이트를 대조하여 분석하세요. 인사이트: {insight}",
            "comp": f"{p_base}[Competitor Analysis] 아래 데이터는 '경쟁사'의 자료입니다. 자사 브랜드({brand_context[:500]}...)와 비교했을 때 이 경쟁사의 취약점과 강점을 분석하세요. 절대 경쟁사를 자사로 착각하지 마세요.",
            "consumer": f"{p_base}[Antithesis] 소비자의 날것의 언어와 결핍(Needs)을 분석하세요. 수집된 데이터의 실질적 페인포인트를 짚으세요.",
            "final": f"{p_base}[Hybrid Strategy]\n1. GAP FINDER: 자사 vs 경쟁사 vs 소비자의 언어 대조를 통해 우리가 잘하는 점/못하는 점을 분석하세요.\n2. Synthesis: '정반합' 결론을 도출하고 실제 운영에 즉시 투입 가능한 DA 카피 소재를 제안하세요.\n자사 운영 인사이트: {insight}"
        }
        res = client.models.generate_content(model="gemini-3-flash-preview", contents=prompts[target_type] + "\n\n데이터:\n" + content[:12000])
        return res.text
    except Exception as e:
        return f"분석 오류: {e}"

# --- 4. 고퀄리티 PDF 엔진 (에러 및 None 방지) ---

class MasterPDF(FPDF):
    def __init__(self):
        super().__init__()
        f_reg, f_bold = "NanumGothic.ttf", "NanumGothicBold.ttf"
        if os.path.exists(f_reg) and os.path.exists(f_bold):
            self.add_font('NG', '', f_reg)
            self.add_font('NG', 'B', f_bold)
            self.font_family_k = 'NG'
        else:
            self.font_family_k = 'Arial'
        _ = self.set_auto_page_break(auto=True, margin=15)
        _ = self.set_margins(20, 15, 20)

    def header(self):
        if hasattr(self, 'font_family_k') and self.page_no() == 1:
            _ = self.set_font(self.font_family_k, 'B', 20)
            _ = self.set_text_color(0, 51, 102)
            _ = self.cell(0, 20, txt="Hybrid Strategic Gap Analysis Report", ln=True, align='C')
            _ = self.ln(5)

    def write_smart(self, text):
        """글자 잘림 방지를 위해 마크다운 기호를 정제하고 긴 문장을 안전하게 출력합니다."""
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                _ = self.ln(5)
                continue
            _ = self.set_font(self.font_family_k, '', 10.5)
            _ = self.set_text_color(50, 50, 50)
            if line.startswith('##') or line.startswith('1.') or line.startswith('2.'): 
                _ = self.set_font(self.font_family_k, 'B', 13)
                line = line.replace('##', '')
            # 인코딩 에러 및 잘림 방지를 위해 특수기호 제거
            clean = re.sub(r'[^\u0000-\u007f\uac00-\ud7af]', '', line.replace('|', ' '))
            _ = self.multi_cell(0, 7, txt=clean)

# --- 5. 단계별 실행 로직 ---

if menu == "STEP 1. 자사 분석 (Thesis)":
    st.title("🏢 STEP 1. 자사 정체성 및 인사이트 분석")
    b_files = st.file_uploader("자사 자료 업로드 (PDF, PPTX, XLSX)", type=["pdf", "pptx", "xlsx"], accept_multiple_files=True)
    b_url = st.text_input("자사 사이트 URL")
    st.session_state['brand_insight'] = st.text_area("💡 실제 운영 인사이트 (예: '앱 기반 부업' 소구의 효율 저조 등)", value=st.session_state['brand_insight'], height=150)
    
    if st.button("자사 분석 실행"):
        if not gemini_key: st.error("Gemini API Key를 입력하세요."); st.stop()
        with st.spinner("자산 및 인사이트 분석 중..."):
            raw_text = extract_text(files=b_files, urls=[b_url])
            st.session_state['brand_analysis'] = analyze_hybrid(raw_text, "brand", st.session_state['brand_insight'])
            _ = st.rerun()
    if st.session_state['brand_analysis']:
        st.markdown(st.session_state['brand_analysis'])

elif menu == "STEP 1.5. 경쟁사 Deep-Dive":
    st.title("⚔️ STEP 1.5. 경쟁사 다중 분석 (최대 3개 세트)")
    st.markdown("자사(원더)와 비교할 경쟁사 브랜드 정보와 자료를 입력하세요.")
    
    c_files = st.file_uploader("경쟁사 전용 자료 업로드", type=["pdf", "pptx", "xlsx"], accept_multiple_files=True)
    
    col1, col2 = st.columns([1, 2])
    with col1:
        c1_n = st.text_input("경쟁사 1 이름")
        c2_n = st.text_input("경쟁사 2 이름")
        c3_n = st.text_input("경쟁사 3 이름")
    with col2:
        c1_u = st.text_input("경쟁사 1 URL")
        c2_u = st.text_input("경쟁사 2 URL")
        c3_u = st.text_input("경쟁사 3 URL")
    
    if st.button("경쟁사 분석 시작"):
        if not st.session_state['brand_analysis']:
            st.error("STEP 1 자사 분석을 먼저 완료하세요."); st.stop()
        with st.spinner("경쟁사 데이터를 수집 및 자사와 대조 중..."):
            all_comp_raw = extract_text(files=c_files, urls=[c1_u, c2_u, c3_u])
            for name in [c1_n, c2_n, c3_n]:
                if name and name.strip():
                    res = requests.post("https://google.serper.dev/search", 
                                        headers={'X-API-KEY': serper_key}, 
                                        json={"q": f"{name} 브랜드 소구점 마케팅 전략", "gl": "kr", "hl": "ko"}).json()
                    if 'organic' in res:
                        all_comp_raw += f"\n[{name} 검색 결과]\n" + "\n".join([r.get('snippet', '') for r in res['organic']])
            
            # [핵심] 자사 분석 결과를 컨텍스트로 주입하여 주객전도 방지
            st.session_state['comp_analysis'] = analyze_hybrid(all_comp_raw, "comp", brand_context=st.session_state['brand_analysis'])
            _ = st.rerun()
    if st.session_state['comp_analysis']:
        st.markdown(st.session_state['comp_analysis'])

elif menu == "STEP 2. 소비자 분석 (Antithesis)":
    st.title("👥 STEP 2. 소비자 리얼 데이터 분석")
    kw = st.text_input("분석 키워드 (예: 제품명 단점)")
    ex = st.text_input("제외 키워드", value="항공, 일본")
    
    if st.button("데이터 수집 시작"):
        if not serper_key: st.error("Serper API Key를 입력하세요."); st.stop()
        with st.spinner("대량 수집 및 스팸 필터링 중..."):
            all_consumer_data = []
            for qs in ["후기", "실망", "단점"]:
                res = requests.post("https://google.serper.dev/search", 
                                    headers={'X-API-KEY': serper_key}, 
                                    json={"q": f"{kw} {qs} -{ex}", "num": 20, "gl": "kr", "hl": "ko"}).json()
                if 'organic' in res:
                    all_consumer_data.extend([{'title': r.get('title'), 'body': r.get('snippet'), 'source': r.get('link')} for r in res['organic']])
            
            st.session_state['consumer_data'] = all_consumer_data
            st.session_state['consumer_analysis'] = analyze_hybrid(str(all_consumer_data), "consumer")
            _ = st.rerun()
    if st.session_state['consumer_analysis']:
        st.markdown(st.session_state['consumer_analysis'])
        _ = st.divider()
        st.subheader("🔍 수집 원본 매핑 (Fact Check)")
        _ = st.dataframe(pd.DataFrame(st.session_state['consumer_data']), use_container_width=True)

elif menu == "STEP 3. 하이브리드 전략 및 PDF":
    st.title("🧠 STEP 3. GAP FINDER & 정반합 통합 전략")
    if st.button("🚀 최종 통합 리포트 생성"):
        if not all([st.session_state['brand_analysis'], st.session_state['comp_analysis'], st.session_state['consumer_analysis']]):
            st.error("이전 단계 분석을 모두 완료하세요.")
        else:
            with st.spinner("데이터 합성 및 전략 수립 중..."):
                combined = f"자사:{st.session_state['brand_analysis']}\n경쟁사:{st.session_state['comp_analysis']}\n소비자:{st.session_state['consumer_analysis']}"
                st.session_state['final_report'] = analyze_hybrid(combined, "final", st.session_state['brand_insight'])
                _ = st.rerun()
    
    if st.session_state['final_report']:
        st.markdown(st.session_state['final_report'])
        _ = st.divider()
        st.subheader("📥 리포트 다운로드")
        
        # [수정] 구문 에러 원천 차단: 개별 if 블록 사용
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            i1 = st.checkbox("자사 포함", value=True)
        with c2:
            i2 = st.checkbox("경쟁사 포함", value=True)
        with c3:
            i3 = st.checkbox("소비자 포함", value=True)
        with c4:
            i4 = st.checkbox("최종전략 포함", value=True)
        
        export_list = []
        if i1: _ = export_list.append(("BRAND ANALYSIS & INSIGHT", st.session_state['brand_analysis']))
        if i2: _ = export_list.append(("COMPETITOR ANALYSIS MATRIX", st.session_state['comp_analysis']))
        if i3: _ = export_list.append(("CONSUMER REAL VOICE ANALYSIS", st.session_state['consumer_analysis']))
        if i4: _ = export_list.append(("FINAL HYBRID STRATEGY REPORT", st.session_state['final_report']))
        
        if export_list:
            pdf = MasterPDF()
            _ = pdf.add_page()
            for title, body in export_list:
                if body:
                    _ = pdf.set_fill_color(240, 240, 240)
                    _ = pdf.set_font(pdf.font_family_k, 'B', 14)
                    _ = pdf.cell(0, 12, txt=f" {title}", ln=True, fill=True)
                    _ = pdf.ln(3)
                    _ = pdf.write_smart(body)
                    _ = pdf.ln(8)
            
            _ = st.download_button(
                label="📥 통합 리포트 PDF 다운로드 (One-Click)",
                data=bytes(pdf.output()),
                file_name="Strategic_Hybrid_Report.pdf",
                mime="application/pdf"
            )
